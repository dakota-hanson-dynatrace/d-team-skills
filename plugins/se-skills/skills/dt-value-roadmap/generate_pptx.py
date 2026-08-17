#!/usr/bin/env python3
"""
dt-value-roadmap: Dynatrace Value Roadmap deck generator.
Usage: python3 generate_pptx.py --customer NAME --tenant URL --data data.json [--output FILE]
"""
import argparse, json, os, shutil, subprocess, datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SKILL_DIR  = os.path.dirname(os.path.abspath(__file__))
ASSETS     = os.path.join(os.path.dirname(SKILL_DIR), 'dynatrace-pptx-skill', 'assets')
TMPL_FIXED = os.path.join(SKILL_DIR, 'dt_template_fixed.pptx')

NAVY       = RGBColor(0x1A, 0x24, 0x40)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEAL       = RGBColor(0x4A, 0xC2, 0xB3)
GREY       = RGBColor(0x6F, 0x74, 0x7F)
RED        = RGBColor(0xC0, 0x39, 0x2B)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
BLUE       = RGBColor(0x29, 0x80, 0xB9)
PURPLE     = RGBColor(0x5E, 0x29, 0xE5)
MAGENTA    = RGBColor(0xC9, 0x3F, 0xDB)
LIGHT_NAVY = RGBColor(0x17, 0x20, 0x36)

TIER_COLOR = {'Start Here': RED, 'High Value': ORANGE, 'Phase 2': BLUE}
TIER_BG    = {
    'Start Here': RGBColor(0x2A, 0x14, 0x14),
    'High Value': RGBColor(0x2A, 0x1E, 0x10),
    'Phase 2':    RGBColor(0x10, 0x1E, 0x2A),
}

# ── OPPORTUNITY LIBRARY ───────────────────────────────────────────────────────
# Each entry: trigger(data)->bool, priority, and format strings keyed {variable}.
# All format strings are rendered with data.format(**data) at build time.
OPPORTUNITIES = [
    {
        'key': 'alert_tuning',
        'trigger': lambda d: d.get('problems_per_day', 0) > 200,
        'priority': 'Start Here',
        'slide_title': 'Alerting Infrastructure Is Built; Now Tune It',
        'card_title':  '{alerting_profiles} alerting profiles are live and routing to teams - {problems_per_day:,}/day volume is the signal to tune',
        'next_bullet': 'Tune alerting - {alerting_profiles} profiles active, {problems_per_day:,}/day volume needs thresholds and AutomationEngine',
        'body': (
            '{customer} has built solid alerting coverage: {alerting_profiles} profiles organized by business unit, '
            'wired to email and Teams channel integrations, and mostly enabled. Davis AI is detecting real problems '
            'and teams are getting notified.\n\n'
            'At {problems_per_day:,} notifications per day, even critical outages risk getting buried. '
            'At this volume, either thresholds are too sensitive or the estate is genuinely unstable - '
            'either way, a critical alert risks being lost in the noise.'
        ),
        'fix': (
            'Tune thresholds and add frequency filters to suppress transient spikes. '
            'Enable disabled Synthetic notifications. '
            'Add AutomationEngine workflows on top of existing notifications: '
            'auto-ticket on AVAILABILITY/ERROR, auto-close on self-resolve within 5 min. '
            'Target: actionable daily volume under 50.'
        ),
        'roadmap_action': 'Tune alerting thresholds + enable AutomationEngine triage workflows',
        'roadmap_value':  'Reduces {problems_per_day:,} daily alerts to an actionable, prioritized volume',
    },
    {
        'key': 'log_monitoring',
        'trigger': lambda d: d.get('log_records_24h', 0) == 0 and d.get('hosts', 0) > 0,
        'priority': 'Start Here',
        'slide_title': 'Log Monitoring Is One Setting Away',
        'card_title':  '{hosts} fully instrumented hosts are ready to stream logs - OneAgent just needs log monitoring enabled',
        'next_bullet': 'Enable log monitoring - OneAgent is already deployed on all {hosts} hosts, one setting to flip',
        'body': (
            'OneAgent is already running on every host and capturing traces. '
            'The same agent can forward system event logs, access logs, and application logs to Dynatrace Grail '
            'with no new infrastructure and no additional deployment.\n\n'
            'Today, when a trace shows an error or slowdown, engineers must manually search log files. '
            'Once enabled, log context appears inline alongside traces - the same DQL query '
            'that shows a slow span also shows the log lines that explain why.\n\n'
            'This is the highest-leverage single action available: one setting, all {hosts} hosts, immediate RCA impact.'
        ),
        'fix': (
            'Enable OneAgent log monitoring in Settings > Log Monitoring. '
            'Configure log sources for the relevant channels (system event logs, web server access logs, application logs). '
            'No new infrastructure required - all logs become queryable alongside traces via DQL.'
        ),
        'roadmap_action': 'Enable OneAgent log monitoring across all {hosts} hosts',
        'roadmap_value':  'Closes RCA blind spot; enables log - trace correlation across all instrumented hosts',
    },
    {
        'key': 'cloud_extension',
        'trigger': lambda d: (
            not d.get('cloud_azure_connected') and
            not d.get('cloud_aws_connected') and
            not d.get('cloud_gcp_connected') and
            d.get('cloud_workloads_exist', False)
        ),
        'priority': 'Start Here',
        'slide_title': 'Extend Observability to {cloud_providers_in_use} as Workloads Move to Cloud',
        'card_title':  'The same OneAgent covering on-prem hosts extends to {cloud_providers_in_use} VMs with no changes to tooling',
        'next_bullet': 'Extend to {cloud_providers_in_use} - connect the subscription before cloud footprint grows further',
        'body': (
            '{customer} runs a hybrid estate and is actively moving workloads to {cloud_providers_in_use}. '
            'The OneAgent deployment model that already covers every on-prem host works identically '
            'on cloud VMs - same agent, same configuration, same DQL queries, unified topology.\n\n'
            'For workloads that run as PaaS services, Dynatrace\'s native cloud integration pulls metrics '
            'and events directly from the cloud provider\'s monitoring APIs without requiring an agent. '
            'Both paths feed into the same Grail data set alongside on-prem data.\n\n'
            'Getting observability in place before migration completes is significantly easier than '
            'retrofitting after. Each workload that moves to cloud without monitoring creates a new blind spot '
            'in an estate that is otherwise well-covered.'
        ),
        'fix': (
            'Connect the cloud subscription via Settings > Cloud and Virtualization > {cloud_providers_in_use}. '
            'Deploy OneAgent to cloud VMs using the same method as on-prem hosts. '
            'Enable container and Kubernetes monitoring if those workloads are part of the migration. '
            'Establish an environment and workload tagging standard before the cloud footprint grows further.'
        ),
        'roadmap_action': 'Connect {cloud_providers_in_use} subscription; extend OneAgent to cloud VMs and PaaS services',
        'roadmap_value':  'Unified on-prem and cloud topology as workloads migrate to {cloud_providers_in_use}',
    },
    {
        'key': 'slos',
        'trigger': lambda d: d.get('slo_count', 0) == 0,
        'priority': 'High Value',
        'slide_title': '{traces_per_day} Traces/Day Can Become Business SLOs Today',
        'card_title':  'The data to power SLOs is already flowing - reliability targets just need to be defined',
        'next_bullet': 'Define SLOs - {traces_per_day} traces/day and {services:,} services give you the data, just need targets set',
        'body': (
            'Dynatrace is capturing every transaction across {services:,} services, '
            '{traces_per_day} per day. That is the exact dataset SLOs are built on. No new instrumentation is needed.\n\n'
            'Right now the team cannot answer "are we meeting our commitments?" with data, and leadership '
            'has no report showing whether applications are performing to target. With SLOs defined, '
            'Dynatrace generates automatic burn-down alerts and status board reporting.\n\n'
            'SLOs can be live within hours. The foundation is already there.'
        ),
        'fix': (
            'Start with: (1) Availability SLO on key production applications (target: 99.9%), '
            '(2) Performance SLO on critical web request services (p95 < 500ms), '
            '(3) Error rate SLO on top 10 services by call volume. '
            'Expand to the full service catalog using management zones.'
        ),
        'roadmap_action': 'Define SLOs on key production applications and top services by call volume',
        'roadmap_value':  'Business-aligned reliability reporting and error budgets across {services:,} services',
    },
    {
        'key': 'workflows',
        'trigger': lambda d: d.get('workflow_count', 0) == 0,
        'priority': 'High Value',
        'slide_title': 'Add AutomationEngine on Top of Active Notifications',
        'card_title':  'Classic notifications are reaching the right teams - AutomationEngine adds the layer that acts on them',
        'next_bullet': 'Enable AutomationEngine workflows on top of existing notification routing',
        'body': (
            'Email and Teams routing is active. Teams are getting notified. That is the hard part - '
            'and {customer} already has it.\n\n'
            'AutomationEngine is the next layer: instead of an engineer reading the email, '
            'opening Dynatrace, creating a ticket, and assigning it by hand - '
            'a workflow does that automatically. On-call engineers shift from mechanical triage to '
            'solving the problems that actually need human judgment.\n\n'
            'Note: AutomationEngine workflows require the AutomationEngine managed IAM policy '
            'bound to the admin group in Account Management.'
        ),
        'fix': (
            'Quick wins (live within one day): '
            '(1) Open ServiceNow or Jira ticket on any AVAILABILITY or ERROR problem, '
            '(2) Teams or Slack alert for problems affecting key production applications, '
            '(3) Auto-close problems that self-resolve within 5 minutes.'
        ),
        'roadmap_action': 'Create workflows: ticket creation, on-call routing, auto-close',
        'roadmap_value':  'Eliminates manual incident triage burden for on-call engineers',
    },
    {
        'key': 'database_extensions',
        'trigger': lambda d: d.get('db_service_pct', 0) > 40,
        'priority': 'High Value',
        'slide_title': 'Unlock Query-Level Visibility Across {db_services} Database Services',
        'card_title':  '{db_services} database services are already detected - extensions add the query depth that traces alone cannot',
        'next_bullet': 'Enable database extensions for {db_services} detected DB services - slow query and pool saturation visibility',
        'body': (
            '{db_services} of {services:,} services are database services - {db_service_pct}% of the service catalog. '
            'Dynatrace already knows they exist, tracks their call counts, and maps them in the topology. '
            'That discovery is the hard part.\n\n'
            'Database extensions build on top of that foundation to add:\n\n'
            '  ▸  Slow query detection with actual query text\n'
            '  ▸  Connection pool saturation alerts before they cascade\n'
            '  ▸  Wait statistics and index usage data\n\n'
            'Database performance is often the primary driver of end-to-end response time variability. '
            'This unlocks the next tier of root cause analysis.'
        ),
        'fix': (
            'Enable database extensions for the specific engines in use - SQL Server, Oracle, MySQL, PostgreSQL, '
            'and MongoDB extensions are all available. '
            'Cross-reference detected database services against known DB hostnames to identify which engines to prioritize first.'
        ),
        'roadmap_action': 'Enable database extensions for DB engines in use',
        'roadmap_value':  'Query-level visibility into the {db_service_pct}% database service tier ({db_services} services)',
    },
    {
        'key': 'rum_enablement',
        'trigger': lambda d: d.get('rum_apps_classic', 0) > 0 and not d.get('grail_rum_active', False),
        'priority': 'Phase 2',
        'slide_title': 'Enable New RUM Experience Across {rum_apps_classic} Configured Applications',
        'card_title':  '{rum_apps_classic} apps have classic RUM and JS beacons in place - the new Grail RUM experience is not yet active',
        'next_bullet': 'Enable new Grail RUM experience for {rum_apps_classic} apps - beacons deployed, platform not yet active',
        'body': (
            '{customer} has already done the instrumentation work: {rum_apps_classic} applications have classic RUM configured '
            'with JS beacons deployed. The agents are in place. What is not yet active is the new Grail RUM '
            'experience, which surfaces session replay, funnel analysis, and full DQL access to user event data.\n\n'
            'Zero user sessions are currently flowing to the new platform - user.events and user.sessions both '
            'return no records. Classic RUM data exists in isolation, separate from the rest of the Grail data set '
            'and not queryable alongside traces and logs.\n\n'
            'This is a configuration step, not a re-instrumentation. The JS beacon infrastructure is already deployed.'
        ),
        'fix': (
            'Enable the new RUM experience for each app in Applications & Microservices > Frontend. '
            'The JS beacon infrastructure is already deployed - this is a configuration step, not re-instrumentation. '
            'Review synthetic monitors for duplicates and QA or production mixing at the same time.'
        ),
        'roadmap_action': 'Enable new Grail RUM experience for {rum_apps_classic} configured apps; audit synthetic monitors',
        'roadmap_value':  'Full user session and funnel visibility queryable alongside traces and logs',
    },
    {
        'key': 'synthetic_cleanup',
        'trigger': lambda d: d.get('synthetic_monitors', 0) > 5 and d.get('rum_apps_classic', 0) == 0,
        'priority': 'Phase 2',
        'slide_title': 'Review {synthetic_monitors} Synthetic Monitors for Gaps and Duplicates',
        'card_title':  '{synthetic_monitors} monitors are active - an audit typically finds QA mixed with prod and coverage gaps on key flows',
        'next_bullet': 'Clean up {synthetic_monitors} synthetic monitors - remove duplicates, separate QA from prod',
        'body': (
            '{customer} has {synthetic_monitors} synthetic monitors running. At this scale, '
            'without a naming convention and environment separation strategy, monitors tend to accumulate duplicates '
            'and mix QA and production measurements into the same alerting baseline.\n\n'
            'Common patterns to look for:\n\n'
            '  ▸  Same endpoint monitored from multiple monitors with no clear ownership\n'
            '  ▸  QA environment monitors sharing alerting profiles with production monitors\n'
            '  ▸  Business-critical user flows with no synthetic coverage\n\n'
            'A one-hour audit typically results in fewer, better-targeted monitors that produce more reliable alerts.'
        ),
        'fix': (
            'Audit all {synthetic_monitors} monitors: (1) Delete or consolidate duplicates covering the same endpoint. '
            '(2) Separate QA monitors into a dedicated management zone with separate alerting thresholds. '
            '(3) Identify business-critical user flows with no synthetic coverage and create clickpath monitors for them.'
        ),
        'roadmap_action': 'Audit and consolidate {synthetic_monitors} synthetic monitors; separate QA from prod',
        'roadmap_value':  'Cleaner alerting baseline; synthetic coverage on critical user flows',
    },
]

TIER_ORDER = ['Start Here', 'High Value', 'Phase 2']

# ── HELPERS ───────────────────────────────────────────────────────────────────
def asset(name):
    return os.path.join(ASSETS, name)

def _ensure_template():
    if os.path.exists(TMPL_FIXED):
        return
    potx = asset('DT_template.potx')
    tmp  = os.path.join(SKILL_DIR, '_tmp_tmpl')
    os.makedirs(tmp, exist_ok=True)
    subprocess.run(['unzip', '-q', potx, '-d', tmp], check=True)
    ct_path = os.path.join(tmp, '[Content_Types].xml')
    ct = open(ct_path).read().replace(
        'presentationml.template.main+xml',
        'presentationml.presentation.main+xml')
    open(ct_path, 'w').write(ct)
    subprocess.run(['bash', '-c', f'cd "{tmp}" && zip -qr "{TMPL_FIXED}" .'], check=True)
    shutil.rmtree(tmp)

def bg(slide, image='dt_content_bg.png'):
    pic = slide.shapes.add_picture(asset(image), 0, 0,
                                   width=Emu(12192000), height=Emu(6858000))
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def title_accent(slide, y=Inches(1.35)):
    bar = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

def add_text(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, font='DT Flow'):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name   = font
    return txb

def add_badge(slide, label, color, x, y):
    box = slide.shapes.add_shape(1, x, y, Inches(1.3), Inches(0.28))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(13); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = 'DT Flow'

def slide_title(slide, text):
    add_text(slide, text, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
             size=26, bold=True, color=WHITE, font='DT Flow')
    title_accent(slide)

def stat_card(slide, big, label, x, y):
    w, h = Inches(2.7), Inches(1.6)
    box  = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_NAVY
    box.line.color.rgb = TEAL; box.line.width = Pt(1.5)
    add_text(slide, big,   x, y + Inches(0.15), w, Inches(0.85),
             size=42, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font='DT Flow')
    add_text(slide, label, x, y + Inches(0.95), w, Inches(0.55),
             size=15, color=GREY, align=PP_ALIGN.CENTER)

def gap_card(slide, title, body, badge_label, badge_color, fix,
             x=Inches(0.5), y=Inches(1.55)):
    add_badge(slide, badge_label, badge_color, x, y)
    add_text(slide, title, x + Inches(1.45), y, Inches(10.8), Inches(0.65),
             size=19, bold=True, color=WHITE, font='DT Flow')
    add_text(slide, body,  x + Inches(0.3),  y + Inches(0.75), Inches(11.6), Inches(2.8),
             size=17, color=WHITE)
    fix_box = slide.shapes.add_shape(1, x, y + Inches(3.8), Inches(12.3), Inches(1.4))
    fix_box.fill.solid(); fix_box.fill.fore_color.rgb = LIGHT_NAVY
    fix_box.line.color.rgb = TEAL; fix_box.line.width = Pt(1)
    add_text(slide, 'Your Next Step',
             x + Inches(0.15), y + Inches(3.85), Inches(2.4), Inches(0.35),
             size=15, bold=True, color=TEAL)
    add_text(slide, fix,
             x + Inches(0.15), y + Inches(4.23), Inches(12.0), Inches(0.9),
             size=15, color=GREY)

def roadmap_table(slide, rows):
    headers = ['#', 'Next Step', 'Priority', 'Value Unlocked']
    cols    = [0.35, 5.8, 1.2, 4.95]
    col_x   = [Inches(0.5)]
    for c in cols[:-1]:
        col_x.append(col_x[-1] + Inches(c))
    row_h = Inches(0.65)
    y     = Inches(1.55)
    for hdr, cw, cx in zip(headers, cols, col_x):
        box = slide.shapes.add_shape(1, cx, y, Inches(cw), row_h)
        box.fill.solid(); box.fill.fore_color.rgb = TEAL; box.line.fill.background()
        tf = box.text_frame; tf.margin_left = Pt(6); tf.margin_top = Pt(4)
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = hdr
        run.font.size = Pt(15); run.font.bold = True
        run.font.color.rgb = NAVY; run.font.name = 'DT Flow'
    for r, (num, action, sev, impact) in enumerate(rows):
        ry  = y + row_h * (r + 1)
        rbg = TIER_BG.get(sev, LIGHT_NAVY)
        for i, (val, cw, cx) in enumerate(zip([num, action, sev, impact], cols, col_x)):
            box = slide.shapes.add_shape(1, cx, ry, Inches(cw), row_h)
            box.fill.solid(); box.fill.fore_color.rgb = rbg
            box.line.color.rgb = GREY; box.line.width = Pt(0.5)
            tf = box.text_frame; tf.margin_left = Pt(6); tf.margin_top = Pt(4)
            tf.word_wrap = True
            p  = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if i in (0, 2) else PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            run.font.size = Pt(14)
            run.font.color.rgb = TIER_COLOR.get(sev, WHITE) if i == 2 else WHITE
            run.font.bold = (i == 2); run.font.name = 'DT Flow'

# ── SCORING + BUILD ───────────────────────────────────────────────────────────
def score(data):
    active = [o for o in OPPORTUNITIES if o['trigger'](data)]
    return sorted(active, key=lambda o: TIER_ORDER.index(o['priority']))

def fmt(s, data):
    try:
        return s.format(**data)
    except (KeyError, ValueError):
        return s

def _auto_stat_cards(data):
    cards = [
        (str(data.get('hosts', 0)),                   'Hosts\nFully Instrumented'),
        (f"{data.get('services', 0):,}",              'Services\nAuto-Discovered'),
        (str(data.get('traces_per_day', '?')),        'Distributed Traces\nPer Day'),
    ]
    if data.get('rum_apps_classic', 0) > 0 and not data.get('grail_rum_active', False):
        cards.append((str(data['rum_apps_classic']), 'Classic RUM Apps\n(Not Yet on Grail)'))
    elif data.get('synthetic_monitors', 0) > 0:
        cards.append((str(data['synthetic_monitors']), 'Synthetic Monitors\nConfigured'))
    else:
        cards.append((f"{data.get('problems_per_week', 0):,}", 'Problems\nLast 7 Days'))
    return cards

def build(data, customer, tenant, output):
    _ensure_template()
    data['customer'] = customer

    prs   = Presentation(TMPL_FIXED)
    TODAY = datetime.date.today().strftime('%B %d, %Y')

    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    active = score(data)
    n      = len(active)

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[0])
    bg(s, 'dt_cover_bg.png')
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = f'{customer}\nDynatrace Value Roadmap'
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE; run.font.size = Pt(42)
                    run.font.bold = True; run.font.name = 'DT Flow'
        elif ph.placeholder_format.idx == 1:
            ph.text = f'Account Review  ·  {TODAY}'
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = TEAL; run.font.size = Pt(22)
                    run.font.name = 'DT Flow'

    # Foundation divider
    s = prs.slides.add_slide(prs.slide_layouts[20])
    bg(s, 'dt_content_bg.png')
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = 'Your Foundation'
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE; run.font.size = Pt(40)
                    run.font.bold = True; run.font.name = 'DT Flow'
    add_text(s, f'What {customer} has already built on Dynatrace - and the value it is delivering today.',
             Inches(0.5), Inches(4.0), Inches(8), Inches(0.6),
             size=20, color=TEAL, italic=True)

    # Stats slide
    s = prs.slides.add_slide(prs.slide_layouts[61])
    bg(s, 'dt_content_bg.png')
    slide_title(s, 'Current State at a Glance')
    cards = data.get('stat_cards') or _auto_stat_cards(data)
    for (big, label), x in zip(cards[:4], [Inches(0.5), Inches(3.35), Inches(6.2), Inches(9.05)]):
        stat_card(s, big, label, x, Inches(1.7))
    add_text(s, 'WHERE YOU GO NEXT', Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.35),
             size=15, bold=True, color=TEAL)
    for i, opp in enumerate(active[:6]):
        add_text(s, f'→  {fmt(opp["next_bullet"], data)}',
                 Inches(0.5), Inches(4.05) + Inches(0.42) * i,
                 Inches(12.3), Inches(0.38), size=16, color=TIER_COLOR[opp['priority']])

    # Next Steps divider
    s = prs.slides.add_slide(prs.slide_layouts[20])
    bg(s, 'dt_content_bg.png')
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = 'Your Next Steps'
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = WHITE; run.font.size = Pt(40)
                    run.font.bold = True; run.font.name = 'DT Flow'
    n_word = {1:'one',2:'two',3:'three',4:'four',5:'five',6:'six',7:'seven',8:'eight'}.get(n, str(n))
    add_text(s, f'{n_word.capitalize()} opportunities to unlock more value from the platform you have already built.',
             Inches(0.5), Inches(4.0), Inches(8), Inches(0.6),
             size=20, color=TEAL, italic=True)

    # Opportunity slides
    for i, opp in enumerate(active, 1):
        s = prs.slides.add_slide(prs.slide_layouts[61])
        bg(s, 'dt_content_bg.png')
        slide_title(s, fmt(f'Opportunity {i} of {n}  - {opp["slide_title"]}', data))
        gap_card(s,
            title=fmt(opp['card_title'], data),
            body=fmt(opp['body'], data),
            badge_label=opp['priority'].upper(),
            badge_color=TIER_COLOR[opp['priority']],
            fix=fmt(opp['fix'], data),
        )

    # Roadmap
    s = prs.slides.add_slide(prs.slide_layouts[61])
    bg(s, 'dt_content_bg.png')
    slide_title(s, 'Your Value Roadmap')
    roadmap_table(s, [
        (str(i), fmt(opp['roadmap_action'], data), opp['priority'], fmt(opp['roadmap_value'], data))
        for i, opp in enumerate(active, 1)
    ])

    # Thank you
    s = prs.slides.add_slide(prs.slide_layouts[63])
    bg(s, 'dt_cover_bg.png')
    add_text(s, 'Questions?', Inches(0.5), Inches(2.3), Inches(7), Inches(1.2),
             size=48, bold=True, color=WHITE, font='DT Flow')
    add_text(s, f'Data sourced via live DQL queries  ·  {tenant}  ·  {TODAY}',
             Inches(0.5), Inches(5.8), Inches(9), Inches(0.5),
             size=15, color=GREY, italic=True)
    grad_box = s.shapes.add_shape(1, Inches(0.5), Inches(3.6), Inches(4.0), Pt(4))
    fill = grad_box.fill; fill.gradient(); fill.gradient_angle = 0
    stops = fill.gradient_stops
    stops[0].color.rgb = PURPLE; stops[1].color.rgb = MAGENTA
    grad_box.line.fill.background()

    prs.save(output)
    print(f'Saved: {output}  ({n} opportunities, {n + 5} slides total)')


if __name__ == '__main__':
    ap = argparse.ArgumentParser(description='Dynatrace Value Roadmap generator')
    ap.add_argument('--customer', required=True, help='Customer display name')
    ap.add_argument('--tenant',   required=True, help='Tenant URL, e.g. abc123.apps.dynatrace.com')
    ap.add_argument('--data',     required=True, help='Path to data.json')
    ap.add_argument('--output',   default=None,  help='Output .pptx path')
    args = ap.parse_args()
    with open(args.data) as f:
        data = json.load(f)
    out = args.output or f'{args.customer.replace(" ", "_")}_Value_Roadmap.pptx'
    build(data, args.customer, args.tenant, out)
